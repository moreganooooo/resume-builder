import os
import shutil
import sys
import tempfile
import unittest
from unittest.mock import patch

SCRIPTS_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "scripts"
)
sys.path.insert(0, SCRIPTS_DIR)

import bootstrap_extractors  # noqa: E402


class TestDetectFileKind(unittest.TestCase):

    def test_pdf(self):
        self.assertEqual(bootstrap_extractors.detect_file_kind("resume.pdf"), "pdf")

    def test_images(self):
        for ext in ("png", "jpg", "jpeg", "heic", "webp"):
            with self.subTest(ext=ext):
                self.assertEqual(
                    bootstrap_extractors.detect_file_kind(f"screenshot.{ext}"), "image"
                )

    def test_docx(self):
        self.assertEqual(bootstrap_extractors.detect_file_kind("resume.docx"), "docx")

    def test_legacy_doc(self):
        self.assertEqual(bootstrap_extractors.detect_file_kind("old_resume.doc"), "doc")

    def test_odt(self):
        self.assertEqual(bootstrap_extractors.detect_file_kind("notes.odt"), "odt")

    def test_pptx(self):
        self.assertEqual(bootstrap_extractors.detect_file_kind("deck.pptx"), "pptx")

    def test_spreadsheet(self):
        self.assertEqual(
            bootstrap_extractors.detect_file_kind("achievements.xlsx"), "spreadsheet"
        )
        self.assertEqual(
            bootstrap_extractors.detect_file_kind("achievements.csv"), "spreadsheet"
        )

    def test_text(self):
        self.assertEqual(bootstrap_extractors.detect_file_kind("notes.txt"), "text")
        self.assertEqual(bootstrap_extractors.detect_file_kind("notes.md"), "text")

    def test_unsupported(self):
        self.assertEqual(
            bootstrap_extractors.detect_file_kind("archive.zip"), "unsupported"
        )

    def test_case_insensitive(self):
        self.assertEqual(bootstrap_extractors.detect_file_kind("RESUME.PDF"), "pdf")


class TestExtractLocalTextDocx(unittest.TestCase):

    def test_round_trip(self):
        import docx

        tmp_dir = tempfile.mkdtemp()
        path = os.path.join(tmp_dir, "fixture.docx")
        doc = docx.Document()
        doc.add_paragraph("Led a 5-person team to launch the Q3 campaign.")
        doc.add_paragraph("Reduced churn by 12% over six months.")
        doc.save(path)

        text = bootstrap_extractors.extract_local_text(path, "docx")

        self.assertIn("Led a 5-person team to launch the Q3 campaign.", text)
        self.assertIn("Reduced churn by 12% over six months.", text)
        shutil.rmtree(tmp_dir)


class TestExtractLocalTextPptx(unittest.TestCase):

    def test_round_trip(self):
        from pptx import Presentation

        tmp_dir = tempfile.mkdtemp()
        path = os.path.join(tmp_dir, "fixture.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Achievements"
        slide.placeholders[1].text = "Grew pipeline by $3M in one quarter."
        prs.save(path)

        text = bootstrap_extractors.extract_local_text(path, "pptx")

        self.assertIn("Grew pipeline by $3M in one quarter.", text)
        shutil.rmtree(tmp_dir)


class TestExtractLocalTextOdt(unittest.TestCase):

    def test_round_trip(self):
        from odf.opendocument import OpenDocumentText
        from odf.text import P

        tmp_dir = tempfile.mkdtemp()
        path = os.path.join(tmp_dir, "fixture.odt")
        doc = OpenDocumentText()
        doc.text.addElement(P(text="Migrated 2,900 accounts into Salesforce."))
        doc.save(path)

        text = bootstrap_extractors.extract_local_text(path, "odt")

        self.assertIn("Migrated 2,900 accounts into Salesforce.", text)
        shutil.rmtree(tmp_dir)


class TestExtractLocalTextSpreadsheet(unittest.TestCase):

    def test_csv_round_trip(self):
        tmp_dir = tempfile.mkdtemp()
        path = os.path.join(tmp_dir, "fixture.csv")
        with open(path, "w", encoding="utf-8") as f:
            f.write("Achievement,Year\nLaunched product X,2022\n")

        text = bootstrap_extractors.extract_local_text(path, "spreadsheet")

        self.assertIn("Launched product X", text)
        self.assertIn("2022", text)
        shutil.rmtree(tmp_dir)

    def test_xlsx_round_trip(self):
        import pandas as pd

        tmp_dir = tempfile.mkdtemp()
        path = os.path.join(tmp_dir, "fixture.xlsx")
        pd.DataFrame({"Achievement": ["Cut costs by 20%"], "Year": [2021]}).to_excel(
            path, index=False
        )

        text = bootstrap_extractors.extract_local_text(path, "spreadsheet")

        self.assertIn("Cut costs by 20%", text)
        shutil.rmtree(tmp_dir)


class TestExtractLocalTextPlain(unittest.TestCase):

    def test_txt_round_trip(self):
        tmp_dir = tempfile.mkdtemp()
        path = os.path.join(tmp_dir, "fixture.txt")
        with open(path, "w", encoding="utf-8") as f:
            f.write("Built a customer onboarding flow from scratch.")

        text = bootstrap_extractors.extract_local_text(path, "text")

        self.assertEqual(text, "Built a customer onboarding flow from scratch.")
        shutil.rmtree(tmp_dir)


class TestExtractLocalTextRejectsUnsupportedKinds(unittest.TestCase):

    def test_pdf_raises(self):
        with self.assertRaises(ValueError):
            bootstrap_extractors.extract_local_text("whatever.pdf", "pdf")

    def test_image_raises(self):
        with self.assertRaises(ValueError):
            bootstrap_extractors.extract_local_text("whatever.png", "image")


class TestConvertLegacyDocToPdf(unittest.TestCase):

    @patch("bootstrap_extractors.shutil.which", return_value=None)
    def test_returns_none_when_libreoffice_unavailable(self, mock_which):
        result = bootstrap_extractors.convert_legacy_doc_to_pdf("old.doc")
        self.assertIsNone(result)


if __name__ == "__main__":
    unittest.main()
