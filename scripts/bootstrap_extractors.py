import os
import shutil
import subprocess

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)

_EXTENSION_KIND_MAP = {
    "pdf": "pdf",
    "png": "image",
    "jpg": "image",
    "jpeg": "image",
    "heic": "image",
    "webp": "image",
    "docx": "docx",
    "doc": "doc",
    "odt": "odt",
    "pptx": "pptx",
    "xlsx": "spreadsheet",
    "xls": "spreadsheet",
    "csv": "spreadsheet",
    "txt": "text",
    "md": "text",
}


def detect_file_kind(path: str) -> str:
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    return _EXTENSION_KIND_MAP.get(ext, "unsupported")


def convert_legacy_doc_to_pdf(path: str) -> str | None:
    soffice = shutil.which("soffice")
    if not soffice:
        return None
    out_dir = os.path.dirname(os.path.abspath(path))
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, path],
        check=True,
        capture_output=True,
    )
    base = os.path.splitext(os.path.basename(path))[0]
    return os.path.join(out_dir, f"{base}.pdf")


def _extract_docx_text(path: str) -> str:
    import docx
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx_text(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _extract_odt_text(path: str) -> str:
    from odf.opendocument import load
    from odf import text as odf_text
    from odf import teletype
    doc = load(path)
    paragraphs = doc.getElementsByType(odf_text.P)
    return "\n".join(teletype.extractText(p) for p in paragraphs)


def _extract_spreadsheet_text(path: str) -> str:
    import pandas as pd
    ext = os.path.splitext(path)[1].lower()
    if ext == ".csv":
        df = pd.read_csv(path)
    else:
        df = pd.read_excel(path)
    return df.to_csv(index=False)


def _extract_plain_text(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def extract_local_text(path: str, kind: str) -> str:
    if kind == "docx":
        return _extract_docx_text(path)
    if kind == "pptx":
        return _extract_pptx_text(path)
    if kind == "odt":
        return _extract_odt_text(path)
    if kind == "spreadsheet":
        return _extract_spreadsheet_text(path)
    if kind == "text":
        return _extract_plain_text(path)
    raise ValueError(f"extract_local_text does not support kind={kind!r} (path={path})")


import sys
from typing import Literal, Optional

from google import genai
from google.genai import types
from pydantic import BaseModel, Field

if SCRIPT_DIR not in sys.path:
    sys.path.insert(0, SCRIPT_DIR)

from gemini_client import GeminiClient  # noqa: E402

EXTRACTION_MODEL = "gemini-3.1-flash-lite"
UPLOAD_MODEL = "gemma-4-31b-it"


class RawAchievement(BaseModel):
    raw_text: str = Field(description="The achievement as written or lightly rephrased for clarity -- never invent detail not in the source.")
    company_hint: Optional[str] = Field(default=None, description="Company/employer name if stated or strongly implied nearby in the source text.")
    date_hint: Optional[str] = Field(default=None, description="Any date or date range mentioned near this achievement, verbatim as written.")
    title_hint: Optional[str] = Field(default=None, description="Job title mentioned near this achievement, if any.")
    confidence: Literal["high", "medium", "low"] = Field(description="high: company/date clearly stated. medium: implied but not explicit. low: no attribution context at all.")


class RawAchievementList(BaseModel):
    achievements: list[RawAchievement]


class Certificate(BaseModel):
    name: str
    issuer: Optional[str] = None
    date: Optional[str] = None


class WorkExperienceEntry(BaseModel):
    company: str
    title: Optional[str] = None
    start_date: Optional[str] = None
    end_date: Optional[str] = None
    achievements: list[str] = Field(default_factory=list)


class ResumeExtraction(BaseModel):
    experience: list[WorkExperienceEntry]
    certifications: list[Certificate] = Field(default_factory=list)


class DocumentClassification(BaseModel):
    doc_type: Literal["resume", "linkedin_export", "recommendation_letter", "achievement_notes", "certificate", "other"]


class IngestionAPIError(RuntimeError):
    """Raised by the ingestion-path extraction functions below (classify_
    document_type/extract_achievements/extract_certificate/extract_resume_
    timeline_and_achievements, text= branch only) when GeminiClient.generate()
    itself failed -- returned raw=None -- rather than genuinely finding
    nothing to extract. raw is None only on a real call failure (permanent
    HTTP error, exhausted retries, malformed response); a successful call
    with nothing to say still returns a string. Without this distinction, a
    failed call (e.g. a 403 from a missing/bad API key) silently produced an
    empty-but-"done" result -- bootstrap_bullet_bank.run_ingestion() catches
    this specifically to checkpoint "failed" (retried next run) instead of
    "done" (permanently skipped). See B16."""


_BASE_EXTRACTION_RULES = """
You are extracting real career achievements from a personal document so they can
become resume bullet points later. Follow these rules strictly:
- Extract only what the source text actually supports. Light rephrasing for
  clarity is fine (fixing grammar, tightening a run-on sentence). Inventing or
  inferring a metric, scope, team size, or outcome that is not stated or
  clearly implied in the source text is NOT fine.
- It is fine to lightly connect obvious dots within this one document (e.g. a
  job title mentioned in one line applying to an achievement described two
  lines later).
- Skip generic filler (objective statements, contact info, soft-skill lists
  with no evidence) -- only extract concrete achievements or accomplishments.
- For each achievement, capture whatever company, date, or job-title context
  appears near it in the source, even if incomplete. Do not guess a company
  name that never appears in the text.
- Set confidence to "high" only when the company AND a date or clear time
  period are both explicitly stated near the achievement. Use "medium" when
  attribution is implied but not explicit. Use "low" when there is no
  attribution context at all.
"""

_EXTRACTION_PROMPTS = {
    "recommendation_letter": _BASE_EXTRACTION_RULES + """
This document is a letter of recommendation written ABOUT this person by
someone else, in third person. Extract only the specific, concrete
achievements or projects the letter describes this person doing -- not the
letter-writer's own opinions, adjectives, or general praise with no
underlying specific action attached.
""",
    "achievement_notes": _BASE_EXTRACTION_RULES + """
This document is free-form notes the person wrote about their own past
achievements. Extract each distinct achievement as its own entry.
""",
    "other": _BASE_EXTRACTION_RULES + """
This document's type is unclear. Extract any concrete, achievement-shaped
statements you find; skip anything that is not a specific accomplishment.
""",
}

_RESUME_EXTRACTION_PROMPT = _BASE_EXTRACTION_RULES + """
This document is a resume or LinkedIn profile export. For each job/role
listed, extract the company name, job title, start/end dates as written,
and every achievement bullet under that role, verbatim or lightly
rephrased for clarity only.

Separately, if this document also lists any certifications or credentials
(e.g. a "Certifications" section), extract those into the certifications
list instead of treating them as achievement bullets -- a credential isn't
an achievement. Do not invent an issuer or date if the document doesn't
state one; use null instead.
"""

_CERTIFICATE_PROMPT = """
You are extracting a professional certificate or credential from a document.
Return the credential's name, issuing organization (if stated), and the date
issued or earned (if stated). Do not invent any of these fields if they are
not present in the source -- use null instead.
"""

_CLASSIFY_PROMPT = """
Classify this document into exactly one category: resume, linkedin_export,
recommendation_letter, achievement_notes, certificate, or other. Use the
filename and the content sample provided.
"""

_FILENAME_HEURISTICS = [
    (("linkedin",), "linkedin_export"),
    (("resume", "cv"), "resume"),
    (("certificate", "certification"), "certificate"),
    (("recommendation", "reference letter", "letter of rec"), "recommendation_letter"),
]


def _api_key() -> str | None:
    return os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")


def _generate_from_upload(path: str, system_prompt: str, response_schema) -> str | None:
    """Uploads a PDF/image file directly to Gemini and returns raw response
    text. GeminiClient's REST client has no file-upload support, so this
    uses the google-genai SDK client directly -- the same proven pattern
    ingest.py already uses for its single-resume parse."""
    client = genai.Client(api_key=_api_key())
    uploaded = client.files.upload(file=path)
    config = types.GenerateContentConfig(
        system_instruction=system_prompt,
        response_mime_type="application/json",
        response_schema=response_schema,
        temperature=0.0,
    )
    response = client.models.generate_content(
        model=UPLOAD_MODEL, contents=[uploaded, "Extract the requested information."], config=config,
    )
    return response.text


def classify_document_type(filename: str, text: str | None, dry_run: bool = False) -> str:
    """Classifies a document by filename heuristic first; falls back to an
    LLM call over its text only when a heuristic doesn't match AND text is
    available. PDFs/images (text=None) with no filename match default to
    'achievement_notes' rather than spending a second multimodal API call
    just to classify -- most real filenames aren't that ambiguous, and this
    keeps ingestion cost proportionate."""
    lowered = filename.lower()
    for keywords, doc_type in _FILENAME_HEURISTICS:
        if any(kw in lowered for kw in keywords):
            return doc_type

    if text is None:
        return "achievement_notes"

    if dry_run:
        print(f"[DRY RUN] would classify {filename!r} via LLM over its text sample.")
        return "other"

    sample = text[:2000]
    raw, _ = GeminiClient.generate(
        model=EXTRACTION_MODEL,
        system_instruction=_CLASSIFY_PROMPT,
        contents=f"Filename: {filename}\n\nContent sample:\n{sample}",
        response_schema=DocumentClassification,
        temperature=0.0,
    )
    if raw is None:
        raise IngestionAPIError(
            f"Gemini API call failed while classifying {filename!r} -- see the WARNING above for the status code."
        )
    data = GeminiClient.parse_json(raw)
    return data.get("doc_type", "other")


def extract_achievements(
    doc_type: str, *, text: str | None = None, upload_path: str | None = None, dry_run: bool = False,
) -> list[RawAchievement]:
    """Extracts achievement-shaped content for doc_type in
    ('recommendation_letter', 'achievement_notes', 'other'). Exactly one of
    text or upload_path must be set."""
    if (text is None) == (upload_path is None):
        raise ValueError("extract_achievements requires exactly one of text or upload_path")

    system_prompt = _EXTRACTION_PROMPTS.get(doc_type, _EXTRACTION_PROMPTS["other"])

    if dry_run:
        print(f"[DRY RUN] would extract achievements (doc_type={doc_type!r}) with prompt:\n{system_prompt}")
        return []

    if upload_path is not None:
        raw = _generate_from_upload(upload_path, system_prompt, RawAchievementList)
    else:
        raw, _ = GeminiClient.generate(
            model=EXTRACTION_MODEL, system_instruction=system_prompt,
            contents=text, response_schema=RawAchievementList, temperature=0.0,
        )
        if raw is None:
            raise IngestionAPIError(
                f"Gemini API call failed while extracting achievements (doc_type={doc_type!r}) -- "
                "see the WARNING above for the status code."
            )
    data = GeminiClient.parse_json(raw) if isinstance(raw, str) else (raw or {})
    return [RawAchievement(**a) for a in data.get("achievements", [])]


def extract_certificate(
    *, text: str | None = None, upload_path: str | None = None, dry_run: bool = False,
) -> Certificate | None:
    """Extracts a single credential (name/issuer/date) from a document
    classified as 'certificate'. Returns None if no credential name was
    found. Exactly one of text or upload_path must be set."""
    if (text is None) == (upload_path is None):
        raise ValueError("extract_certificate requires exactly one of text or upload_path")

    if dry_run:
        print("[DRY RUN] would extract a certificate.")
        return None

    if upload_path is not None:
        raw = _generate_from_upload(upload_path, _CERTIFICATE_PROMPT, Certificate)
    else:
        raw, _ = GeminiClient.generate(
            model=EXTRACTION_MODEL, system_instruction=_CERTIFICATE_PROMPT,
            contents=text, response_schema=Certificate, temperature=0.0,
        )
        if raw is None:
            raise IngestionAPIError(
                "Gemini API call failed while extracting a certificate -- see the WARNING above for the status code."
            )
    data = GeminiClient.parse_json(raw) if isinstance(raw, str) else (raw or {})
    if not data.get("name"):
        return None
    return Certificate(**data)


def extract_resume_timeline_and_achievements(
    *, text: str | None = None, upload_path: str | None = None, dry_run: bool = False,
) -> ResumeExtraction:
    """Used specifically for documents classified as 'resume' or
    'linkedin_export' -- extracts the employment timeline (company, title,
    dates) AND each role's listed achievement bullets in one pass, since
    the source document already states which company each bullet belongs
    to (no separate attribution step needed for these). Also extracts any
    embedded certifications/credentials section into
    ResumeExtraction.certifications, rather than treating a credential as
    an achievement bullet. Exactly one of text or upload_path must be set."""
    if (text is None) == (upload_path is None):
        raise ValueError("extract_resume_timeline_and_achievements requires exactly one of text or upload_path")

    if dry_run:
        print("[DRY RUN] would extract resume/LinkedIn timeline, achievements, and certifications.")
        return ResumeExtraction(experience=[], certifications=[])

    if upload_path is not None:
        raw = _generate_from_upload(upload_path, _RESUME_EXTRACTION_PROMPT, ResumeExtraction)
    else:
        raw, _ = GeminiClient.generate(
            model=EXTRACTION_MODEL, system_instruction=_RESUME_EXTRACTION_PROMPT,
            contents=text, response_schema=ResumeExtraction, temperature=0.0,
        )
        if raw is None:
            raise IngestionAPIError(
                "Gemini API call failed while extracting the resume/LinkedIn timeline -- "
                "see the WARNING above for the status code."
            )
    data = GeminiClient.parse_json(raw) if isinstance(raw, str) else (raw or {})
    return ResumeExtraction(**data) if data else ResumeExtraction(experience=[], certifications=[])


class ContactInfo(BaseModel):
    full_name: Optional[str] = None
    email: Optional[str] = None
    phone: Optional[str] = None
    location: Optional[str] = None
    linkedin_url: Optional[str] = None
    portfolio_url: Optional[str] = None


class RecommendationQuote(BaseModel):
    name: Optional[str] = None
    title: Optional[str] = None
    quote: Optional[str] = None


class RoleSuggestions(BaseModel):
    secondary_roles: list[str] = Field(default_factory=list)


class TagDefinition(BaseModel):
    name: str
    persona_description: str
    keywords: list[str] = Field(default_factory=list)


class TagTaxonomy(BaseModel):
    tags: list[TagDefinition] = Field(default_factory=list)


class LedgerEntry(BaseModel):
    label: str
    value: str
    employer: str = ""


class NamedLedgerItem(BaseModel):
    name: str
    employer: str = ""


class LedgerExtraction(BaseModel):
    metrics: list[LedgerEntry] = Field(default_factory=list)
    tools: list[NamedLedgerItem] = Field(default_factory=list)
    projects: list[NamedLedgerItem] = Field(default_factory=list)


_CONTACT_INFO_PROMPT = """
Extract this person's contact/identity information from the document text:
full name, email, phone, location (city/state), LinkedIn profile URL, and
portfolio/personal website URL if present. Use null for any field not
explicitly present in the text -- never invent a value.
"""

_RECOMMENDATION_QUOTE_PROMPT = """
This document is a letter of recommendation written about this person by
someone else. Extract the single strongest, most specific endorsing quote
from the letter, along with the name and title of the person who wrote it,
if stated. Use null for any field not explicitly present -- never invent
a name, title, or quote that isn't really in the letter.
"""

_SECONDARY_ROLES_PROMPT = """
Given a person's primary target job titles and a sample of their real
achievements, suggest 2-3 adjacent job titles they could also reasonably
target -- roles that draw on the same underlying skills but aren't
identical to the primary titles. Only suggest real, standard job titles;
do not invent a title that wouldn't make sense to a recruiter.
"""

_BACKGROUND_GUIDE_PROMPT = """
You are drafting a short narrative career-background guide from this
person's own resume/LinkedIn summary, recommendation-letter excerpts, and
achievement notes. Write 2-4 short paragraphs in third person describing
how their background/skills came together and what makes their combination
of experience distinctive. Only describe things directly supported by the
source text -- do not invent employers, dates, or accomplishments not
present in the source. If the source material is too thin to say anything
specific, write a short, honest, general paragraph instead of padding it
with invented detail.
"""

_TAG_TAXONOMY_PROMPT = """
Given a candidate's target job titles and a sample of their real
achievement bullets, define a tag taxonomy for categorizing their resume
bullets by skill/theme area -- 4-8 short, single-word, lowercase tags
(e.g. "email", "ops", "leadership", "design") that meaningfully partition
the kinds of work actually evident in their target roles and achievements.
Base the taxonomy entirely on the target roles and achievement text given
-- do not assume a specific industry (e.g. marketing) unless the input
evidence actually supports it; a different candidate's evidence should
produce a different taxonomy.

For each tag, provide:
- name: the tag word itself, no brackets (e.g. "email", not "[email]")
- persona_description: one sentence naming the job archetypes/role types
  this tag's content is most relevant to (e.g. "email marketing, lifecycle
  marketing, or CRM/ESP campaign roles")
- keywords: 5-12 lowercase keywords/phrases that would plausibly appear in
  a bullet belonging to this tag -- prefer specific terms over generic ones

Always include exactly one catch-all tag named "generalist" with
persona_description "general or cross-functional roles" and an EMPTY
keywords list (empty keywords is how the rest of the system recognizes a
catch-all tag that matches everything, not a mistake to fix).
"""

_VOICE_ANCHORS_PROMPT = """
You are drafting a "voice anchors" file from this person's own writing --
real answers to interview-style questions, cover letter excerpts, or
any other first-person writing sample text they've provided. The goal is
to capture genuine turns of phrase and themes worth echoing later, not to
summarize their career (a separate background guide already does that).

Output 3-6 entries in this exact markdown shape, one per real
question/topic actually evident in the source text:

### <a short question or topic, phrased the way an interviewer or cover
letter prompt might ask it -- e.g. "Why would you be a good fit for X"
or "Tell me about a time you had to prioritize under pressure">

<one or two sentences paraphrasing their real answer/point, third person>

> <OPTIONAL -- only if the source text contains an actual quotable
line worth echoing verbatim; a real sentence or near-verbatim close
paraphrase, in first person, in quotes. Omit this line entirely for an
entry with no standout phrase to pull -- do not invent one.>

Only draw from what's actually in the source text -- do not invent
questions, answers, or quotes that aren't grounded in it. If the source
material is too thin to produce even 3 genuine entries, output fewer
rather than padding with invented ones.
"""

_LEDGER_PROMPT = """
Given a list of resume achievement bullets, each prefixed with the employer
it belongs to in brackets like "[Acme Corp] Grew pipeline by 22%", extract
three things:
- metrics: every quantified result mentioned, as a (label, value) pair
  where value is the number/stat exactly as written (e.g. "22% reply rate")
- tools: every named tool, platform, or piece of software mentioned
- projects: every named project or initiative mentioned

For every metric, tool, and project, also set "employer" to the bracketed
company name of the bullet it came from. If the exact same tool or project
name appears under more than one employer's bullets, emit one entry per
employer rather than merging them. Leave employer as an empty string only
if a bullet has no bracketed company at all -- never guess or invent one.

Only include things explicitly stated in the text. Do not invent numbers,
tool names, or project names that aren't there.
"""


def extract_contact_info(
    *, text: str | None = None, upload_path: str | None = None, dry_run: bool = False,
) -> ContactInfo:
    """Extracts identity/contact fields from a resume or LinkedIn export's
    text. Exactly one of text or upload_path must be set."""
    if (text is None) == (upload_path is None):
        raise ValueError("extract_contact_info requires exactly one of text or upload_path")

    if dry_run:
        print("[DRY RUN] would extract contact info.")
        return ContactInfo()

    if upload_path is not None:
        raw = _generate_from_upload(upload_path, _CONTACT_INFO_PROMPT, ContactInfo)
    else:
        raw, _ = GeminiClient.generate(
            model=EXTRACTION_MODEL, system_instruction=_CONTACT_INFO_PROMPT,
            contents=text, response_schema=ContactInfo, temperature=0.0,
        )
    data = GeminiClient.parse_json(raw) if isinstance(raw, str) else (raw or {})
    return ContactInfo(**data) if data else ContactInfo()


def extract_recommendation_quote(
    *, text: str | None = None, upload_path: str | None = None, dry_run: bool = False,
) -> RecommendationQuote | None:
    """Extracts a real quote + attribution from a recommendation letter.
    Returns None if no usable quote was found. Exactly one of text or
    upload_path must be set."""
    if (text is None) == (upload_path is None):
        raise ValueError("extract_recommendation_quote requires exactly one of text or upload_path")

    if dry_run:
        print("[DRY RUN] would extract a recommendation quote.")
        return None

    if upload_path is not None:
        raw = _generate_from_upload(upload_path, _RECOMMENDATION_QUOTE_PROMPT, RecommendationQuote)
    else:
        raw, _ = GeminiClient.generate(
            model=EXTRACTION_MODEL, system_instruction=_RECOMMENDATION_QUOTE_PROMPT,
            contents=text, response_schema=RecommendationQuote, temperature=0.0,
        )
    data = GeminiClient.parse_json(raw) if isinstance(raw, str) else (raw or {})
    if not data.get("quote"):
        return None
    return RecommendationQuote(**data)


def suggest_secondary_roles(
    primary_roles: list[str], achievements_text: str, dry_run: bool = False,
) -> list[str]:
    """Suggests 2-3 adjacent target job titles based on confirmed primary
    roles and a sample of real achievement text."""
    if dry_run:
        print("[DRY RUN] would suggest secondary target roles.")
        return []

    raw, _ = GeminiClient.generate(
        model=EXTRACTION_MODEL, system_instruction=_SECONDARY_ROLES_PROMPT,
        contents=f"Primary roles: {', '.join(primary_roles)}\n\nAchievements:\n{achievements_text[:4000]}",
        response_schema=RoleSuggestions, temperature=0.0,
    )
    data = GeminiClient.parse_json(raw)
    return data.get("secondary_roles", [])


def generate_tag_taxonomy(
    primary_roles: list[str], secondary_roles: list[str], achievements_text: str, dry_run: bool = False,
) -> TagTaxonomy:
    """Derives a per-profile bullet-tag taxonomy (name/persona_description/
    keywords per tag) from this candidate's own target roles and real
    achievement text -- replaces what used to be Morgan's hardcoded
    marketing-specific tags (see profile_paths.tags()'s docstring), so
    auto-tagging and claim/context filtering work for any field, not just
    marketing/sales. Falls back to a single "generalist" tag with no
    keywords if the model returns nothing usable, so callers always get at
    least one tag to work with."""
    if dry_run:
        print("[DRY RUN] would generate a tag taxonomy.")
        return TagTaxonomy()

    roles_text = ", ".join(primary_roles + secondary_roles)
    raw, _ = GeminiClient.generate(
        model=EXTRACTION_MODEL, system_instruction=_TAG_TAXONOMY_PROMPT,
        contents=f"Target roles: {roles_text}\n\nSample achievements:\n{achievements_text[:6000]}",
        response_schema=TagTaxonomy, temperature=0.2,
    )
    data = GeminiClient.parse_json(raw)
    taxonomy = TagTaxonomy(**data) if data else TagTaxonomy()
    if not taxonomy.tags:
        taxonomy = TagTaxonomy(tags=[
            TagDefinition(name="generalist", persona_description="general or cross-functional roles", keywords=[]),
        ])
    return taxonomy


def draft_background_guide(source_texts: list[str], dry_run: bool = False) -> str:
    """Synthesizes a short narrative background guide from resume/rec-letter/
    achievement-notes text already gathered during bootstrap ingestion."""
    if dry_run:
        print("[DRY RUN] would draft a background guide.")
        return ""

    combined = "\n\n---\n\n".join(t for t in source_texts if t)[:8000]
    raw, _ = GeminiClient.generate(
        model=EXTRACTION_MODEL, system_instruction=_BACKGROUND_GUIDE_PROMPT,
        contents=combined, temperature=0.4,
    )
    return (raw or "").strip()


def draft_voice_anchors(source_texts: list[str], dry_run: bool = False) -> str:
    """Drafts voice-anchors.md -- real turns of phrase and themes worth
    echoing later, pulled from whatever first-person writing-sample text
    got ingested during bootstrap (cover letters, interview-prep answers,
    etc.). Degrades to an empty string (not an invented example) when
    there's nothing usable to draw from; every consumer of voice-anchors.md
    (orchestrator.py, rewrite_bullets.py) already checks os.path.exists()
    first, so an empty/missing file is a fully supported "no signal yet"
    state, not an error."""
    if dry_run:
        print("[DRY RUN] would draft voice anchors.")
        return ""

    combined = "\n\n---\n\n".join(t for t in source_texts if t)[:8000]
    if not combined:
        return ""
    raw, _ = GeminiClient.generate(
        model=EXTRACTION_MODEL, system_instruction=_VOICE_ANCHORS_PROMPT,
        contents=combined, temperature=0.4,
    )
    return (raw or "").strip()


def extract_ledger_entries(achievements_text: str, dry_run: bool = False) -> LedgerExtraction:
    """Derives simple metrics/tools/projects lists from already-extracted
    achievement bullets, for the verified_metrics/tools/projects.json
    ledger files. Single model call over achievements_text as given --
    callers with a large profile should chunk first via
    extract_ledger_entries_chunked() rather than passing raw text longer
    than one call can safely see (see that function's docstring)."""
    if dry_run:
        print("[DRY RUN] would extract ledger entries (metrics/tools/projects).")
        return LedgerExtraction()

    raw, _ = GeminiClient.generate(
        model=EXTRACTION_MODEL, system_instruction=_LEDGER_PROMPT,
        contents=achievements_text, response_schema=LedgerExtraction, temperature=0.0,
    )
    data = GeminiClient.parse_json(raw)
    return LedgerExtraction(**data) if data else LedgerExtraction()


LEDGER_CHUNK_CHARS = 6000  # matches this call's previous hardcoded truncation cap


def _chunk_lines(text: str, max_chars: int) -> list[str]:
    """Splits text into newline-bounded chunks each under max_chars --
    never breaks a single line (one bullet) across two chunks."""
    lines = text.split("\n")
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0
    for line in lines:
        added_len = len(line) + (1 if current else 0)
        if current and current_len + added_len > max_chars:
            chunks.append("\n".join(current))
            current, current_len = [], 0
            added_len = len(line)
        current.append(line)
        current_len += added_len
    if current:
        chunks.append("\n".join(current))
    return chunks


def extract_ledger_entries_chunked(achievements_text: str, dry_run: bool = False) -> LedgerExtraction:
    """extract_ledger_entries(), but safe for a profile with more
    achievement text than a single call can see. The old direct call
    silently truncated achievements_text to 6000 chars, meaning a person
    with several companies' worth of bullets could have later companies
    dropped from their verified_metrics/tools/projects.json entirely with
    no warning. This splits on bullet boundaries into <=6000-char chunks,
    extracts each independently, and merges the results -- deduping exact
    repeats (same label/value or name, same employer) that show up in more
    than one chunk, while keeping the same name/label attributed to two
    DIFFERENT employers as two separate entries (per _LEDGER_PROMPT)."""
    if dry_run:
        print("[DRY RUN] would extract ledger entries (metrics/tools/projects).")
        return LedgerExtraction()

    chunks = _chunk_lines(achievements_text, LEDGER_CHUNK_CHARS)
    metrics: list[LedgerEntry] = []
    tools: list[NamedLedgerItem] = []
    projects: list[NamedLedgerItem] = []
    seen_metrics: set = set()
    seen_tools: set = set()
    seen_projects: set = set()

    for chunk in chunks:
        result = extract_ledger_entries(chunk)
        for m in result.metrics:
            key = (m.label.strip().lower(), m.value.strip().lower(), m.employer.strip().lower())
            if key not in seen_metrics:
                seen_metrics.add(key)
                metrics.append(m)
        for t in result.tools:
            key = (t.name.strip().lower(), t.employer.strip().lower())
            if key not in seen_tools:
                seen_tools.add(key)
                tools.append(t)
        for p in result.projects:
            key = (p.name.strip().lower(), p.employer.strip().lower())
            if key not in seen_projects:
                seen_projects.add(key)
                projects.append(p)

    return LedgerExtraction(metrics=metrics, tools=tools, projects=projects)
